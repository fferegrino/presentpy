from pathlib import Path

import click
import nbformat

from presentpy.code_slide_source import CodeSlideSource


class NotebookProcessor:
    def __init__(self, notebook_path: Path, output_device):
        self.notebook_path = notebook_path
        self.notebook = nbformat.read(self.notebook_path, as_version=4)
        self.output_device = output_device

    def process(self):
        for cell in self.notebook.cells:
            if cell.cell_type == "code":
                code_slide_source = CodeSlideSource.from_source_code(cell.source)
                self.output_device.write(code_slide_source)

from pptx import Presentation
class PptxOutputDevice:
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    def __init__(self, output_path: Path):
        self.presentation = Presentation()
    def write(self, code_slide_source: CodeSlideSource):
        title_slide  = self.presentation.slide_layouts[PptxOutputDevice.SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = self.presentation.slides.add_slide(title_slide)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = code_slide_source.title or ""
        content.text = code_slide_source.code

    def save(self, output_path: Path):
        self.presentation.save(output_path)




@click.command()
@click.argument('notebook', type=click.Path(exists=True))
def run_presentpy(notebook):

    notebook_path = Path(notebook)
    output_file_name = Path(notebook_path.stem + ".pptx")

    pptx_output_device = PptxOutputDevice(notebook_path)
    notebook_processor = NotebookProcessor(notebook_path, pptx_output_device)
    notebook_processor.process()

    pptx_output_device.save(output_file_name)


    print(f'Running {notebook_path}...')

if __name__ == '__main__':
    run_presentpy()
