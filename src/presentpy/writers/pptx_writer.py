from pathlib import Path
from typing import Any, Dict

import pkg_resources
import pygments.styles
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Pt
from pygments.token import Token

from presentpy.code_slide_source import CodeSlideSource, get_parsed_lines

CARRIAGE_RETURN = "\x0A"

EXTRA_STYLES = {
    "light": {
        Token.Punctuation: RGBColor(102, 102, 102),
        Token.Literal.String.Single: RGBColor(64, 112, 160),
        Token.Literal.String.Doc: RGBColor(64, 112, 160),
        Token.Literal.Number.Integer: RGBColor(64, 160, 112),
        Token.Keyword.Namespace: RGBColor(0, 112, 32),
        Token.Name.Builtin.Pseudo: RGBColor(27, 82, 167),
        Token.Name.Function.Magic: RGBColor(49, 0, 250),
        Token.Comment.Single: RGBColor(65, 127, 127),
        Token.Keyword.Constant: RGBColor(0, 112, 32),
    },
    "dark": {
        Token.Literal.Number.Integer: RGBColor.from_string("F78C6C"),
        Token.Comment.Single: RGBColor.from_string("546E7A"),
    },
}


def get_theme(theme: str = "light") -> Dict[Any, RGBColor]:
    style_name = "friendly" if theme == "light" else "material"

    style = pygments.styles.get_style_by_name(style_name)
    token_colors = {}
    for token, str_style in style.styles.items():
        if not str_style:
            continue
        _, _, color = str_style.partition("#")
        if not color:
            continue

        pad = 1 if len(color) == 3 else 2
        token_colors[token] = RGBColor(*[int(color[i : i + pad], 16) for i in range(0, len(color), pad)])

    token_colors.update(EXTRA_STYLES.get(theme, {}))

    return token_colors


class PptxWriter:
    TEMPLATE_CODE_SLIDE_INDEX = 2

    def __init__(self, theme: str = "light"):
        template = pkg_resources.resource_filename("presentpy", f"slide_templates/template-{theme}.pptx")
        self.presentation = Presentation(template)
        self.theme = get_theme(theme)

    def _write_code_slide(self, slide_title, parsed_tokens, lines_to_highlight):
        slide_layout = self.presentation.slide_layouts[PptxWriter.TEMPLATE_CODE_SLIDE_INDEX]
        slide = self.presentation.slides.add_slide(slide_layout)

        title = slide.shapes.title
        if slide_title:
            title.text = slide_title

        content = slide.shapes.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.bullet = False

        for line_number, line_content in enumerate(parsed_tokens, 1):
            for kind, text in line_content:
                run = p.add_run()
                run.text = text
                font = run.font
                font.bold = line_number in lines_to_highlight
                font.color.rgb = self.theme.get(kind, RGBColor(0, 0, 0))
                font.size = Pt(14)
            run = p.add_run()
            run.font.size = Pt(14)
            run.text = CARRIAGE_RETURN

    def write(self, code_slide_source: CodeSlideSource):
        parsed_lines = get_parsed_lines(code_slide_source.code)

        # We need to add this -1 so that there is AT LEAST ONE slide rendered
        # The -1 is because no line number will match, therefore the source code will be rendered without highlighting
        highlights = [[-1]]
        highlights.extend(code_slide_source.highlights)

        for hl in highlights:
            self._write_code_slide(code_slide_source.title, parsed_lines, hl)

    def save(self, output_path: Path):
        self.presentation.save(output_path)
